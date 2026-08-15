"""Generate DDW Sales Deck PPT — 10 slides, dark+gold theme."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# === Design Tokens ===
DARK  = RGBColor(0x0A, 0x0A, 0x1A)
GOLD  = RGBColor(0xE8, 0xB8, 0x6D)
PANEL = RGBColor(0x14, 0x14, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BGL   = RGBColor(0xF8, 0xF6, 0xF0)
TD    = RGBColor(0x0F, 0x0B, 0x0A)
TM    = RGBColor(0x6A, 0x6A, 0x80)
TW    = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xEF, 0x44, 0x44)
ORN   = RGBColor(0xF5, 0x9E, 0x0B)
SW = Inches(13.333); SH = Inches(7.5)
FC = 'PingFang SC'; FE = 'Helvetica'

def R(s, x, y, w, h, f=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if f: sh.fill.solid(); sh.fill.fore_color.rgb = f
    else: sh.fill.background()
    sh.line.fill.background(); return sh

def C(s, x, y, w, h, f=WHITE):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.08
    if f: sh.fill.solid(); sh.fill.fore_color.rgb = f
    else: sh.fill.background()
    sh.line.fill.background(); return sh

def T(s, x, y, w, h, t, sz=14, b=False, c=TD, a=PP_ALIGN.LEFT, fn=FC):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, ln in enumerate(t.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = a
        run = p.add_run(); run.text = ln; run.font.name = fn; run.font.size = Pt(sz)
        run.font.bold = b; run.font.color.rgb = c
        rPr = run._r.get_or_add_rPr()
        for ea in rPr.findall(qn('a:ea')): rPr.remove(ea)
        etree.SubElement(rPr, qn('a:ea')).set('typeface', FC)
    return tb

def HDR(s, t, sub=None, pn=None):
    R(s, Inches(0.5), Inches(0.4), Inches(0.10), Inches(0.50), GOLD)
    T(s, Inches(0.7), Inches(0.35), Inches(10), Inches(0.5), t, 26, True, TD)
    if sub: T(s, Inches(0.7), Inches(0.88), Inches(10), Inches(0.35), sub, 11, False, TM)
    if pn: T(s, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3), f'{pn} / 10', 10, False, TM, PP_ALIGN.RIGHT)

def BG(s, c=BGL):
    b = R(s, 0, 0, SW, SH, c)
    sp = b._element.getparent(); sp.remove(b._element); sp.insert(2, b._element)

def FT(s): T(s, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3), 'DDW AI Hub', 9, False, TM)

prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BL = prs.slide_layouts[6]

# S1: Title
s = prs.slides.add_slide(BL); R(s, 0, 0, SW, SH, DARK)
T(s, Inches(0.8), Inches(0.6), Inches(8), Inches(0.4), 'DDW AI HUB', 10, True, GOLD, fn=FE)
T(s, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2), 'DDW AI \xb5\xd7\xd7\xf9\xc6\xbd\xcc\xa8', 64, True, TW)
T(s, Inches(0.8), Inches(3.1), Inches(11), Inches(0.7), '\xb0\xd1\xc6\xf3\xd2\xb5 AI \xc2\xe4\xb5\xd8\xa3\xac\xb4\xd3 6 \xb8\xf6\xd4\xc2\xcb\xf5\xb6\xcc\xb5\xbd 6 \xcc\xec', 28, False, GOLD)
R(s, Inches(0.8), Inches(4.0), Inches(2), Inches(0.06), GOLD)
T(s, Inches(0.8), Inches(4.3), Inches(6), Inches(1.5), '\xc3\xe6\xcf\xf2\xd6\xd0\xd0\xa1\xd6\xc6\xd4\xec\xc6\xf3\xd2\xb5\n\xd2\xbb\xb4\xce\xb2\xbf\xca\xf0\xa3\xac\xd6\xd5\xc9\xed\xca\xdc\xd3\xc3', 16, False, RGBColor(0xAA,0xAA,0xCC))

# S2-S10: Content slides
slides = [
    (2, '\xc6\xf3\xd2\xb5 AI \xc2\xe4\xb5\xd8 5 \xb4\xf3\xcd\xb4\xb5\xe3',
     [('\xb6\xe0\xc4\xa3\xd0\xcd\xbb\xec\xc2\xd2','\xb6\xc8\xb3\xa4\xd3\xc3 GPT, \xb0\xb2\xc8\xab\xd2\xaa DeepSeek, IT \xb5\xa3\xd0\xc4 Ollama','3 \xcc\xd7 API \xb8\xf7\xd7\xd4\xc3\xdc\xd4\xbf => Token \xca\xa7\xbf\xd8',RED),
      ('\xb2\xe5\xbc\xfe\xb0\xf3\xcb\xc0','\xc2\xf2\xc1\xcb AI \xd6\xfa\xca\xd6, \xcf\xeb\xbb\xbb OCR','\xca\xfd\xbe\xdd\xc7\xa8\xc5\xc2\xb7\xfe\xce\xf1\xc6\xf2, 3 \xb8\xf6\xd4\xc2',RED),
      ('\xb1\xbe\xb5\xd8\xb2\xbf\xca\xf0\xb9\xf3','\xb9\xa4\xb3\xa7\xca\xfd\xbe\xdd\xb2\xbb\xc4\xdc\xc9\xcf\xd4\xc6','\xcb\xbd\xd3\xd0\xbb\xaf\xb1\xa8\xbc\xdb 80 \xcd\xf2, \xd6\xd0\xd0\xa1\xc6\xf3\xd2\xb5\xb1\xbb\xc5\xc5\xb3\xfd',ORN),
      ('\xd4\xb1\xb9\xa4\xd3\xc3\xb2\xbb\xc6\xf0\xc0\xb4','\xcd\xc6\xc1\xcb 3 \xb8\xf6\xd4\xc2 AI \xb9\xa4\xbe\xdf','\xc8\xd5\xbb\xee\xb2\xbb\xb5\xbd 8%, \xb1\xc8 Excel \xb6\xbc\xb2\xbb\xc8\xe7',ORN),
      ('ERP / MES \xb8\xee\xc1\xd1','\xd2\xd1\xd3\xd0 ERP / MES / WMS','AI \xb9\xa4\xbe\xdf\xbd\xf8\xb2\xbb\xc8\xa5\xd2\xb2\xb3\xf6\xb2\xbb\xc0\xb4',ORN)]),
    (3, '\xb7\xbd\xb0\xb8: DDW AI \xb5\xd7\xd7\xf9', [('\xbc\xab\xbc\xf2\xb2\xbf\xca\xf0','\xd2\xbb\xcc\xf5\xc3\xfc\xc1\xee\xc6\xf4\xb6\xaf\nddw server start\n\xc8\xfd\xb6\xcb\xcd\xb3\xd2\xbb'), ('\xb2\xe5\xbc\xfe\xc9\xfa\xcc\xac','\xb1\xea\xd7\xbc SDK => plugin marketplace\n6 \xb8\xf6\xcf\xd6\xd3\xd0\xb2\xe5\xbc\xfe\nERP / MES \xbf\xc9\xb6\xa8\xd6\xc6'), ('LLM \xcd\xf8\xb9\xd8','MiniMax M3, DeepSeek V4 Pro\nOllama \xb1\xbe\xb5\xd8\xbd\xb5\xbc\xb6\n\xc8\xfd\xbc\xb6\xd7\xd4\xb6\xaf\xc7\xd0\xbb\xbb'), ('FDE \xbd\xbb\xb8\xb6','\xbe\xad\xc0\xed\xbc\xec\xb2\xe2\xb5\xbd\xbd\xbb\xb8\xb0\n2 \xcc\xec\xc4\xda\xcd\xea\xb3\xc9\n\xb2\xe5\xbc\xfe\xb2\xbb\xb9\xbb => \xcf\xd6\xb3\xa1\xbf\xaa\xb7\xa2')]),
]

for pn, title, items in slides:
    s = prs.slides.add_slide(BL); BG(s); HDR(s, title, None, pn)
    for i, (t, *rest) in enumerate(items):
        x = Inches(0.5 + i * 2.6); y = Inches(1.8)
        desc = rest[0] if len(rest) >= 1 else ''
        sub = rest[1] if len(rest) >= 2 else ''
        clr = rest[2] if len(rest) >= 3 else GOLD
        C(s, x, y, Inches(2.3), Inches(2.8), WHITE)
        R(s, x + Inches(0.1), y + Inches(0.1), Inches(2.1), Inches(0.06), clr)
        T(s, x + Inches(0.15), y + Inches(0.3), Inches(2.0), Inches(0.35), f'{i+1}. {t}', 14, True, TD)
        T(s, x + Inches(0.15), y + Inches(0.75), Inches(2.0), Inches(1.8), '\n'.join(filter(None, [desc, sub])), 10, False, TM)
    FT(s)

# S10: Closing
s = prs.slides.add_slide(BL); R(s, 0, 0, SW, SH, DARK)
T(s, Inches(1), Inches(1.5), Inches(11), Inches(1.0), '\xd2\xbb\xb8\xf6\xc8\xcb\xbe\xcd\xca\xc7\xd2\xbb\xd6\xa7\xcd\xc5\xb6\xd3', 48, True, TW)
R(s, Inches(1), Inches(2.7), Inches(2), Inches(0.06), GOLD)
T(s, Inches(1), Inches(3.0), Inches(11), Inches(1.0), 'DDW AI Hub \xb7 FDE \xc3\xe2\xb7\xd1\xd7\xa4\xb3\xa1 PoC\n\xb0\xd1\xc6\xf3\xd2\xb5 AI \xc2\xe4\xb5\xd8\xa3\xac\xb4\xd3 6 \xb8\xf6\xd4\xc2\xcb\xf5\xb6\xcc\xb5\xbd 6 \xcc\xec', 20, False, GOLD)

out = '/Users/chenye/workspace/ddw-ai-hub/cloud-llm/ddw-ai-hub/docs/ddw_sales_deck.pptx'
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'OK: {out} ({len(prs.slides)} slides, {os.path.getsize(out)} B)')
