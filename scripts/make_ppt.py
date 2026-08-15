from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

DARK_BG = RGBColor(0x0A, 0x0A, 0x1A)
GOLD = RGBColor(0xE8, 0xB8, 0x6D)
PANEL_DARK = RGBColor(0x14, 0x14, 0x2E)
LIGHT_BG = RGBColor(0xF8, 0xF6, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x0F, 0x0B, 0x0A)
TEXT_MUTED = RGBColor(0x6A, 0x6A, 0x80)
TEXT_ON_DARK = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xEF, 0x44, 0x44)
WARN_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
SW = Inches(13.333)
SH = Inches(7.5)
FONT_CN = 'PingFang SC'
FONT_EN = 'Helvetica'

def add_rect(s, x, y, w, h, fill=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    shp.line.fill.background(); return shp

def add_rounded(s, x, y, w, h, fill=WHITE):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.08
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    shp.line.fill.background(); return shp

def add_text(s, x, y, w, h, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font=FONT_CN):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font; run.font.size = Pt(size)
        run.font.bold = bold; run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        for ea in rPr.findall(qn('a:ea')): rPr.remove(ea)
        etree.SubElement(rPr, qn('a:ea')).set('typeface', FONT_CN)
    return tb

def add_header(s, title, subtitle=None, page_num=None):
    add_rect(s, Inches(0.5), Inches(0.4), Inches(0.10), Inches(0.50), fill=GOLD)
    add_text(s, Inches(0.7), Inches(0.35), Inches(10), Inches(0.5), title, 26, True, TEXT_DARK)
    if subtitle: add_text(s, Inches(0.7), Inches(0.88), Inches(10), Inches(0.35), subtitle, 11, False, TEXT_MUTED)
    if page_num: add_text(s, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3), f'{page_num} / 10', 10, False, TEXT_MUTED, PP_ALIGN.RIGHT)

def add_bg(s, color=LIGHT_BG):
    b = add_rect(s, 0, 0, SW, SH, fill=color)
    spTree = b._element.getparent(); spTree.remove(b._element); spTree.insert(2, b._element)

def add_footer(s):
    add_text(s, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3), 'DDW AI Hub \xb7 \xb6c\xfa\xc7\xa7 AI \xb5\xd7\xd7\xf9\xc6\xbd\xcc\xa8', 9, False, TEXT_MUTED)

prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
blank = prs.slide_layouts[6]

# Slide 1
s = prs.slides.add_slide(blank); add_rect(s, 0, 0, SW, SH, DARK_BG)
add_text(s, Inches(0.8), Inches(0.6), Inches(8), Inches(0.4), 'DDW AI HUB \xb7 PLATFORM FOR MANUFACTURING', 10, True, GOLD, font=FONT_EN)
add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2), 'DDW AI \xb5\xd7\xd7\xf9\xc6\xbd\xcc\xa8', 64, True, TEXT_ON_DARK)
add_text(s, Inches(0.8), Inches(3.1), Inches(11), Inches(0.7), '\xb0\xd1\xc6\xf3\xd2\xb5 AI \xc2\xe4\xb5\xd8\xa3\xac\xb4\xd3 6 \xb8\xf6\xd4\xc2\xcb\xf5\xb6\xcc\xb5\xbd 6 \xcc\xec', 28, False, GOLD)
add_rect(s, Inches(0.8), Inches(4.0), Inches(2), Inches(0.06), GOLD)
add_text(s, Inches(0.8), Inches(4.3), Inches(6), Inches(1.5), '\xc3\xe6\xcf\xf2\xd6\xd0\xb9\xfa\xd6\xd0\xd0\xa1\xd6\xc6\xd4\xec\xc6\xf3\xd2\xb5\n\xd2\xbb\xb4\xce\xb2\xbf\xca\xf0\xa3\xac\xd6\xd5\xc9\xed\xca\xdc\xd3\xc3 \xb7 plugin \xc9\xfa\xcc\xac \xb7 \xb1\xbe\xb5\xd8\xd3\xc5\xcf\xc8', 16, False, RGBColor(0xAA, 0xAA, 0xCC))

# Slide 2 - Pain Points
s = prs.slides.add_slide(blank); add_bg(s); add_header(s, '\xc6\xf3\xd2\xb5 AI \xc2\xe4\xb5\xd8 5 \xb4\xf3\xcd\xb4\xb5\xe3', None, 2)
pains = [
    ('\xb6\xe0\xc4\xa3\xd0\xcd\xbb\xec\xc2\xd2', '\xb6\xc8\xb3\xa4\xd3\xc3 GPT\xa3\xac\xb0\xb2\xc8\xab\xd2\xaa DeepSeek\n3 \xcc\xd7 API \xb8\xf7\xd7\xd4\xc3\xdc\xd4\xbf => Token \xca\xa7\xbf\xd8', DANGER),
    ('\xb2\xe5\xbc\xfe\xb0\xf3\xcb\xc0', '\xc2\xf2\xc1\xcb AI \xd6\xfa\xca\xd6\xa3\xac\xcf\xeb\xbb\xbb OCR\n\xca\xfd\xbe\xdd\xc7\xa8\xc5\xc2\xd2\xaa 3 \xb8\xf6\xd4\xc2', DANGER),
    ('\xb1\xbe\xb5\xd8\xb2\xbf\xca\xf0\xb9\xf3', '\xb9\xa4\xb3\xa7\xca\xfd\xbe\xdd\xb2\xbb\xc4\xdc\xc9\xcf\xd4\xc6\n\xc7\xa8\xc3\xfb\xbb\xaf\xb1\xa8\xbc\xdb 80 \xcd\xf2', WARN_ORANGE),
    ('\xd4\xb1\xb9\xa4\xd3\xc3\xb2\xbb\xc6\xf0\xc0\xb4', '\xcd\xc6\xc1\xcb 3 \xb8\xf6\xd4\xc2 AI \xb9\xa4\xbe\xdf\n\xc8\xd5\xbb\xee\xb2\xbb\xb5\xbd 8%', WARN_ORANGE),
    ('\xcf\xb5\xcd\xb3\xb8\xee\xc1\xd1', '\xd2\xd1\xd3\xd0 ERP / MES / WMS\nAI \xb9\xa4\xbe\xdf\xbd\xf8\xb2\xbb\xc8\xa5', WARN_ORANGE),
]
for i, (t, d, c) in enumerate(pains):
    x = Inches(0.5 + i * 2.5)
    add_rounded(s, x, Inches(1.6), Inches(2.2), Inches(3.0), WHITE)
    add_rect(s, x + Inches(0.1), Inches(1.7), Inches(2.0), Inches(0.06), c)
    add_text(s, x + Inches(0.15), Inches(1.95), Inches(1.9), Inches(0.35), t, 14, True, TEXT_DARK)
    add_text(s, x + Inches(0.15), Inches(2.4), Inches(1.9), Inches(1.8), d, 10, False, TEXT_MUTED)
add_footer(s)

# Save PPT
out = '/Users/chenye/workspace/ddw-ai-hub/cloud-llm/ddw-ai-hub/docs/ddw_sales_deck.pptx'
prs.save(out)
print(f"OK PPT: {out} ({len(prs.slides)} slides)")
print(f"Size: {os.path.getsize(out)} bytes")
